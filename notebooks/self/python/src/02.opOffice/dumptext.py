
import os

from docx import Document
from pptx import Presentation

# 提取指定目录下的 word，ppt 中的文字，并保存在 dump.txt 中


def extractWord(file_name):
    # 将word文档的内容一行一行的读取
    docx = Document(file_name)

    text = []

    # 将word文档的内容一行一行的读取
    for paragraph in docx.paragraphs:
        text.append(paragraph.text)

    for table in docx.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)

    return text


def extractPPT(file_name):
    prs = Presentation(file_name)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    return text


# 所有文件的文件名
all_files = []


def find_files(dir_name):
    # 递归查找文件
    allowed_files = ('docx', 'pptx', 'ppt')

    lsdir = os.listdir(dir_name)
    dirs = [i for i in lsdir if os.path.isdir(os.path.join(dir_name, i))]
    if dirs:
        for i in dirs:
            find_files(os.path.join(dir_name, i))
    files = [i for i in lsdir if os.path.isfile(os.path.join(dir_name, i))]
    for f in files:
        if (f.lower().endswith(allowed_files)):
            all_files.append(os.path.join(dir_name, f))


def dump_files():
    # input: all_files, 遍历其中每一个文件
    results = []

    def error(f): return print('no match: %'.format(f))
    fn = error

    for file_name in all_files:
        # 忽略大小为零的文件
        statinfo = os.stat(file_name)
        if (statinfo.st_size <= 0):
            continue

        if file_name.endswith('.docx'):
            fn = extractWord
        if file_name.endswith(('.ppt', 'pptx')):
            fn = extractPPT

        t = fn(file_name)
        results = results + t

        print('parse {}'.format(file_name))

    # 使用固定的文件名
    lines = [s.strip() for s in results if len(s.strip()) > 0]
    with open('dump.txt', 'w') as f:
        for line in lines:
            f.write("%s\n" % line)

    print('dump complete. {} lines'.format(len(lines)))


if __name__ == "__main__":
    dir_name = '/Users/ygzheng/Documents/project/Support/圣象/itsp2/90.Tools/90.ISC/005 访谈及调研/访谈纪要'
    # dir_name = '.'
    find_files(dir_name)

    dump_files()

    print('done.')
