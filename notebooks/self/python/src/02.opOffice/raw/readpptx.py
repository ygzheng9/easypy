from pptx import Presentation

import jieba

from pyecharts import options as opts
from pyecharts.charts import Page, WordCloud
from pyecharts.globals import SymbolType


path_to_presentation = './1.pptx'
prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)


# print(text_runs)


words_list = []


def seperate(line):
    global words_list
    # words = jieba.cut(line, cut_all=False)
    words = jieba.lcut(line, cut_all=False)

    words_list = words_list + words


for line in text_runs:
    seperate(line)

# print(words_list)


count_list = []


def word_count(list):
    myset = set(list)  # myset是另外一个列表，里面的内容是mylist里面的无重复 项
    for item in myset:
        if len(item) == 1:
            continue

        count_list.append({'name': item, 'count': list.count(item)})


word_count(words_list)


sorted_list = sorted(count_list, key=lambda x: x['count'], reverse=True)

print(sorted_list)


words = [(elem['name'], elem['count']) for elem in sorted_list]


def wordcloud_base():
    c = (
        WordCloud()
        .add("", words, word_size_range=[20, 100])
        .set_global_opts(title_opts=opts.TitleOpts(title="报告信息"))
    )
    return c


t = wordcloud_base()
t.render()

print('done.')
